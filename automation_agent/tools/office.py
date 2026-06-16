import shutil
import subprocess
from pathlib import Path

from automation_agent.tools.base import ToolResult


class OfficeTools:
    def __init__(self, safety):
        self.safety = safety

    def execute(self, operation, arguments):
        handler = getattr(self, operation, None)
        if handler is None:
            raise ValueError(f"Unsupported office operation: {operation}")
        return handler(**arguments)

    def create_docx(self, destination, title, paragraphs):
        from docx import Document

        destination_path = self.safety.validate_path(destination)
        destination_path.parent.mkdir(parents=True, exist_ok=True)
        document = Document()
        document.add_heading(title, level=0)
        for paragraph in paragraphs:
            document.add_paragraph(str(paragraph))
        document.save(destination_path)
        return ToolResult(True, f"Created {destination_path.name}.", [str(destination_path)])

    def create_pptx(self, destination, title, slides):
        from pptx import Presentation

        destination_path = self.safety.validate_path(destination)
        destination_path.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        for index, slide_data in enumerate(slides):
            layout = presentation.slide_layouts[0 if index == 0 else 1]
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = slide_data.get("title") or (title if index == 0 else "Section")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "\n".join(slide_data.get("bullets", []))
        presentation.save(destination_path)
        return ToolResult(True, f"Created {destination_path.name}.", [str(destination_path)])

    def create_spreadsheet(self, destination, headers, rows, sheet_name="Report"):
        from openpyxl import Workbook

        destination_path = self.safety.validate_path(destination)
        destination_path.parent.mkdir(parents=True, exist_ok=True)
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = sheet_name[:31]
        sheet.append(list(headers))
        for row in rows:
            sheet.append(list(row))
        sheet.freeze_panes = "A2"
        sheet.auto_filter.ref = sheet.dimensions
        workbook.save(destination_path)
        return ToolResult(True, f"Created {destination_path.name}.", [str(destination_path)])

    def weekly_excel_report(self, source, destination):
        from docx import Document
        from openpyxl import load_workbook

        source_path = self.safety.validate_path(source, must_exist=True)
        destination_path = self.safety.validate_path(destination)
        destination_path.parent.mkdir(parents=True, exist_ok=True)
        workbooks = [source_path] if source_path.is_file() else sorted(source_path.glob("*.xlsx"))
        if not workbooks:
            raise ValueError(f"No XLSX workbooks were found in {source_path}.")

        document = Document()
        document.add_heading("Weekly Spreadsheet Report", level=0)
        document.add_paragraph(f"Generated locally from {len(workbooks)} workbook(s).")
        summaries = []
        for workbook_path in workbooks:
            workbook = load_workbook(workbook_path, read_only=True, data_only=True)
            document.add_heading(workbook_path.name, level=1)
            for sheet in workbook.worksheets:
                populated = sum(
                    1
                    for row in sheet.iter_rows()
                    if any(cell.value not in (None, "") for cell in row)
                )
                summary = {
                    "workbook": workbook_path.name,
                    "sheet": sheet.title,
                    "rows": populated,
                    "columns": sheet.max_column,
                }
                summaries.append(summary)
                document.add_paragraph(
                    f"{sheet.title}: {populated} populated rows, {sheet.max_column} columns.",
                    style="List Bullet",
                )
            workbook.close()
        document.save(destination_path)
        return ToolResult(
            True,
            f"Created weekly report from {len(workbooks)} workbook(s).",
            [str(destination_path)],
            {"sheets": summaries},
        )

    def convert_to_pdf(self, source, destination_dir):
        source_path = self.safety.validate_path(source, must_exist=True)
        output_dir = self.safety.validate_path(destination_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            raise RuntimeError("LibreOffice is not installed or is not available on PATH.")
        completed = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(source_path)],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
        output = output_dir / f"{source_path.stem}.pdf"
        if completed.returncode != 0 or not output.exists():
            raise RuntimeError((completed.stderr or completed.stdout or "LibreOffice conversion failed.").strip())
        return ToolResult(True, f"Converted {source_path.name} to PDF.", [str(output)])
